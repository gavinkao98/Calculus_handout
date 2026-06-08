"""Generate calculus_handout_workflow.pptx in Midnight Canvas dark theme."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- Palette (Midnight Canvas) ----------
BG_DARK   = RGBColor(0x0B, 0x0C, 0x10)
CARD_BG   = RGBColor(0x16, 0x19, 0x21)
CARD_BG2  = RGBColor(0x1F, 0x23, 0x2D)
TITLE     = RGBColor(0x4C, 0xC9, 0xF0)
ACCENT    = RGBColor(0xF9, 0xA8, 0x25)
BODY      = RGBColor(0xE8, 0xE8, 0xF0)
SUBTLE    = RGBColor(0x8A, 0x8A, 0x96)
CORAL     = RGBColor(0xFF, 0x6B, 0x6B)
EMERALD   = RGBColor(0x06, 0xD6, 0xA0)
MATH_BLUE = RGBColor(0x7D, 0xF9, 0xFF)
DIVIDER   = RGBColor(0x32, 0x36, 0x42)

# ---------- Fonts ----------
FONT_TITLE_CN = "Microsoft JhengHei"   # title (Traditional Chinese friendly)
FONT_BODY_CN  = "Microsoft JhengHei"
FONT_MONO     = "Consolas"

# ---------- Embedded screenshots ----------
import os.path
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "embed_images")
IMG_CH4_OPEN = os.path.join(IMG_DIR, "shot_02.png")    # Ch4 chapter opening
IMG_CH1_PAGE = os.path.join(IMG_DIR, "shot_03.png")    # Ch1 strategy + example
IMG_ENV_TBL  = os.path.join(IMG_DIR, "shot_04.png")    # environment cheat sheet
IMG_GITHUB   = os.path.join(IMG_DIR, "shot_07.png")    # GitHub narration.md view

# ---------- Geometry ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_bg(slide, color=BG_DARK):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, color=BODY, bold=False,
             italic=False, font=FONT_BODY_CN, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        runs = [{"text": text}]
    else:
        runs = text
    for i, r in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = r.get("align", align)
        p.line_spacing = line_spacing
        if "space_after" in r:
            p.space_after = Pt(r["space_after"])
        run = p.add_run()
        run.text = r["text"]
        f = run.font
        f.name = r.get("font", font)
        f.size = Pt(r.get("size", size))
        f.bold = r.get("bold", bold)
        f.italic = r.get("italic", italic)
        f.color.rgb = r.get("color", color)
    return box


def add_rect(slide, x, y, w, h, fill=CARD_BG, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, x, y, w, h, fill=CARD_BG, line=None, radius=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_line(slide, x1, y1, x2, y2, color=ACCENT, width=1.5):
    x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.shadow.inherit = False
    return line


def add_arrow(slide, bx, by, ex, ey, color=ACCENT, width=1.5):
    """Arrow connector from (bx,by) to (ex,ey). All coords must be int EMU."""
    bx, by, ex, ey = int(bx), int(by), int(ex), int(ey)
    line = slide.shapes.add_connector(1, bx, by, ex, ey)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    line.shadow.inherit = False
    return line


def add_image_with_frame(slide, path, x, y, w, h):
    """Add an image at (x,y) with size (w,h), surrounded by a thin gold frame
    that doubles as a content separator on the dark canvas."""
    # frame inset 0.04"
    pad = Inches(0.04)
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x - pad, y - pad, w + 2 * pad, h + 2 * pad
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = ACCENT
    frame.line.fill.background()
    frame.shadow.inherit = False
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    pic.shadow.inherit = False
    return pic


def add_footer(slide, idx=None, total=None):
    """Auto-numbered footer: idx defaults to len(prs.slides) so callers don't
    have to track positions across re-orderings."""
    if idx is None:
        idx = len(prs.slides)
    if total is None:
        total = TOTAL
    add_line(slide, Inches(0.5), Inches(7.10), Inches(12.833), Inches(7.10),
             color=DIVIDER, width=0.5)
    add_text(slide, Inches(0.5), Inches(7.18), Inches(8), Inches(0.3),
             [{"text": "微積分講義製作流程", "size": 10, "color": SUBTLE,
               "font": FONT_BODY_CN}],
             align=PP_ALIGN.LEFT)
    add_text(slide, Inches(11), Inches(7.18), Inches(1.833), Inches(0.3),
             [{"text": f"{idx} / {total}", "size": 10, "color": SUBTLE,
               "font": FONT_BODY_CN, "bold": True}],
             align=PP_ALIGN.RIGHT)


def add_slide_title(slide, text, *, accent=ACCENT):
    # small accent square + title
    add_rect(slide, Inches(0.55), Inches(0.55), Inches(0.10), Inches(0.45),
             fill=accent)
    add_text(slide, Inches(0.85), Inches(0.45), Inches(11.5), Inches(0.7),
             [{"text": text, "size": 30, "color": TITLE, "bold": True,
               "font": FONT_TITLE_CN}],
             anchor=MSO_ANCHOR.MIDDLE)


TOTAL = 15


# ===========================================================
# Slide 1 — Cover
# ===========================================================
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    # decorative dot pattern in corner (subtle)
    for r in range(4):
        for c in range(4):
            add_rect(s, Inches(11.6 + c * 0.32), Inches(0.5 + r * 0.32),
                     Inches(0.05), Inches(0.05),
                     fill=DIVIDER)
    # top brand label
    add_text(s, Inches(0.9), Inches(0.7), Inches(8), Inches(0.4),
             [{"text": "CALCULUS HANDOUT  ·  WORKFLOW BRIEFING",
               "size": 12, "color": ACCENT, "bold": True,
               "font": "Calibri"}])
    # main title (centered vertically lower-mid)
    add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.2),
             [{"text": "微積分講義製作流程", "size": 60, "color": TITLE,
               "bold": True, "font": FONT_TITLE_CN}],
             align=PP_ALIGN.LEFT)
    # gold rule
    add_rect(s, Inches(0.9), Inches(3.95), Inches(2.6), Inches(0.05),
             fill=ACCENT)
    # subtitle
    add_text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.6),
             [{"text": "手稿  →  講義  →  Manim 影片",
               "size": 26, "color": BODY, "font": FONT_TITLE_CN}],
             align=PP_ALIGN.LEFT)
    # tagline
    add_text(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(0.5),
             [{"text": "三模式工作流  ·  GitHub 版本控制  ·  雙產出",
               "size": 16, "color": SUBTLE, "font": FONT_TITLE_CN}],
             align=PP_ALIGN.LEFT)
    # date block
    add_rect(s, Inches(0.9), Inches(6.0), Inches(0.06), Inches(0.45),
             fill=ACCENT)
    add_text(s, Inches(1.05), Inches(5.95), Inches(6), Inches(0.55),
             [{"text": "2026-04-30", "size": 18, "color": BODY,
               "bold": True, "font": "Calibri"}])
    # very subtle footer
    add_text(s, Inches(11), Inches(7.18), Inches(1.833), Inches(0.3),
             [{"text": f"1 / {TOTAL}", "size": 10, "color": SUBTLE,
               "font": FONT_BODY_CN, "bold": True}],
             align=PP_ALIGN.RIGHT)


# ===========================================================
# Slide 2 — Overview + GitHub
# ===========================================================
def slide_overview():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "全景：一份手稿，三個產出")

    # three product cards
    cards = [
        ("PDF", "PDF 講義", "單面 A4，高中生自學自足", TITLE),
        ("MP4", "Manim 教學影片", "補強，不是必要", ACCENT),
        ("GIT", "GitHub 版本控制", "PR 管理變更，CI 自動把關", EMERALD),
    ]
    card_y = Inches(1.5)
    card_h = Inches(2.0)
    card_w = Inches(3.95)
    gap = Inches(0.18)
    start_x = Inches(0.55)
    for i, (badge, name, desc, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_round_rect(s, x, card_y, card_w, card_h,
                       fill=CARD_BG, radius=0.05)
        # left accent bar
        add_rect(s, x, card_y, Inches(0.08), card_h, fill=color)
        # badge
        add_text(s, x + Inches(0.4), card_y + Inches(0.3),
                 Inches(1.2), Inches(0.45),
                 [{"text": badge, "size": 14, "color": color,
                   "bold": True, "font": "Calibri"}])
        # name
        add_text(s, x + Inches(0.4), card_y + Inches(0.75),
                 card_w - Inches(0.6), Inches(0.6),
                 [{"text": name, "size": 24, "color": BODY,
                   "bold": True, "font": FONT_TITLE_CN}])
        # desc
        add_text(s, x + Inches(0.4), card_y + Inches(1.4),
                 card_w - Inches(0.6), Inches(0.6),
                 [{"text": desc, "size": 15, "color": SUBTLE,
                   "font": FONT_BODY_CN}],
                 line_spacing=1.3)

    # Lower section: positioning + docs + CI in three columns
    y0 = Inches(4.0)
    sec_h = Inches(2.7)
    col_w = Inches(3.95)
    cols = [
        ("風格定位",
         "Stewart / Rogawski 自學 register",
         "目標讀者：高中生自學進入大學微積分；講義自足，影片補強"),
        ("階層式文件",
         "README · CONTENT_SPEC",
         "CONTENT_ROADMAP · CONTENT_QUICKSTART · CONTENT_EXERCISES"),
        ("CI 五道",
         "book_style_lint · book_preamble_smoketest",
         "book_docs_lint · manim_storyboard_lint · latexmk"),
    ]
    for i, (h, line1, line2) in enumerate(cols):
        x = start_x + i * (col_w + gap)
        add_round_rect(s, x, y0, col_w, sec_h, fill=CARD_BG2, radius=0.04)
        add_text(s, x + Inches(0.35), y0 + Inches(0.3),
                 col_w - Inches(0.5), Inches(0.5),
                 [{"text": h, "size": 16, "color": ACCENT,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, x + Inches(0.35), y0 + Inches(0.85),
                 col_w - Inches(0.5), Inches(0.6),
                 [{"text": line1, "size": 15, "color": BODY,
                   "bold": True, "font": "Calibri"}],
                 line_spacing=1.25)
        add_text(s, x + Inches(0.35), y0 + Inches(1.5),
                 col_w - Inches(0.5), Inches(1.1),
                 [{"text": line2, "size": 13, "color": SUBTLE,
                   "font": "Calibri"}],
                 line_spacing=1.35)

    add_footer(s)


# ===========================================================
# Slide 3 — Manuscript input
# ===========================================================
def slide_manuscript():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "手稿輸入：多人協作的起點")

    bullets = [
        ("01", "多位老師分章撰寫手稿",
         "格式不一：手寫稿 / 既存 LaTeX / 掃描檔皆可"),
        ("02", "Claude 將手稿轉為專案規範的 LaTeX",
         "依 CONTENT_SPEC 的環境集、註記、引用、圖片規則"),
        ("03", "手稿即主軸：不跳過、不重排、不改寫數學內容",
         "原方法、原變數、原證明手法保留；註記式調整收進 roadmap Open questions"),
        ("04", "擴充由 Claude 加上，但每筆都必須打標記",
         "% expansion:<category> 標明來源與動機；review 時可逐筆審視"),
    ]
    y = Inches(1.65)
    h = Inches(1.18)
    gap = Inches(0.18)
    for num, head, sub in bullets:
        # number badge
        add_round_rect(s, Inches(0.55), y, Inches(1.0), h,
                       fill=CARD_BG2, radius=0.10)
        add_text(s, Inches(0.55), y, Inches(1.0), h,
                 [{"text": num, "size": 32, "color": ACCENT,
                   "bold": True, "font": "Calibri"}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # main row
        add_round_rect(s, Inches(1.7), y, Inches(11.08), h,
                       fill=CARD_BG, radius=0.04)
        add_text(s, Inches(2.0), y + Inches(0.18),
                 Inches(10.7), Inches(0.5),
                 [{"text": head, "size": 19, "color": BODY,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, Inches(2.0), y + Inches(0.65),
                 Inches(10.7), Inches(0.45),
                 [{"text": sub, "size": 13.5, "color": SUBTLE,
                   "font": FONT_BODY_CN}], line_spacing=1.3)
        y += h + gap

    add_footer(s)


# ===========================================================
# Slide 4 — Three modes (CORE)
# ===========================================================
def slide_three_modes():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "三模式工作流（核心機制）")

    # State machine row 1: 新手稿 → Mode A → (Mode B) → 簽核
    row1_y = Inches(1.55)
    box_h = Inches(0.85)
    boxes_r1 = [
        ("新手稿", BODY, CARD_BG2, Inches(0.55), Inches(1.5)),
        ("Mode A 起草", TITLE, CARD_BG, Inches(2.55), Inches(2.4)),
        ("Mode B 審核 (可選)", SUBTLE, CARD_BG2, Inches(5.55), Inches(2.6)),
        ("簽核合併到 main", EMERALD, CARD_BG, Inches(8.75), Inches(2.7)),
    ]
    centers_r1 = []
    for label, color, fill, x, w in boxes_r1:
        add_round_rect(s, x, row1_y, w, box_h, fill=fill, radius=0.10)
        add_text(s, x, row1_y, w, box_h,
                 [{"text": label, "size": 16, "color": color,
                   "bold": True, "font": FONT_TITLE_CN}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        centers_r1.append((x, x + w))

    # arrows row 1
    arrow_y = row1_y + box_h // 2
    for i in range(len(boxes_r1) - 1):
        x_start = centers_r1[i][1] + Inches(0.05)
        x_end = centers_r1[i + 1][0] - Inches(0.05)
        add_arrow(s, x_start, arrow_y, x_end, arrow_y,
                  color=ACCENT, width=2)

    # downward branch from "簽核" to row 2 start
    branch_x = (centers_r1[3][0] + centers_r1[3][1]) // 2
    branch_y_top = row1_y + box_h + Inches(0.05)
    branch_y_bot = branch_y_top + Inches(0.55)
    add_arrow(s, branch_x, branch_y_top, branch_x, branch_y_bot,
              color=CORAL, width=2)
    # Place the label to the LEFT of the arrow so it doesn't sit on the line.
    add_text(s, branch_x - Inches(2.7), branch_y_top + Inches(0.15),
             Inches(2.5), Inches(0.3),
             [{"text": "簽核後若需加深  ↘", "size": 11, "color": CORAL,
               "italic": True, "font": FONT_BODY_CN}],
             align=PP_ALIGN.RIGHT)

    # Row 2: Mode C → Mode B (required)
    row2_y = Inches(3.2)
    boxes_r2 = [
        ("Mode C 加深", ACCENT, CARD_BG, Inches(8.75), Inches(2.7)),
        ("Mode B 審核 (必要)", CORAL, CARD_BG2, Inches(5.55), Inches(2.6)),
    ]
    # right box first (Mode C)
    label, color, fill, x, w = boxes_r2[0]
    add_round_rect(s, x, row2_y, w, box_h, fill=fill, radius=0.10,
                   line=ACCENT)
    add_text(s, x, row2_y, w, box_h,
             [{"text": label, "size": 16, "color": color,
               "bold": True, "font": FONT_TITLE_CN}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrow leftward
    label2, color2, fill2, x2, w2 = boxes_r2[1]
    add_round_rect(s, x2, row2_y, w2, box_h, fill=fill2, radius=0.10,
                   line=CORAL)
    add_text(s, x2, row2_y, w2, box_h,
             [{"text": label2, "size": 16, "color": color2,
               "bold": True, "font": FONT_TITLE_CN}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrow Mode C → Mode B (right to left)
    arrow_y2 = row2_y + box_h // 2
    add_arrow(s, x - Inches(0.05), arrow_y2,
              x2 + w2 + Inches(0.05), arrow_y2,
              color=ACCENT, width=2)

    # Three mode descriptions across bottom
    desc_y = Inches(4.6)
    desc_h = Inches(2.25)
    col_w = Inches(3.95)
    gap = Inches(0.18)
    descs = [
        ("Mode A — 起草", TITLE,
         "手稿轉 LaTeX；圍繞手稿擴充。",
         "每筆擴充加標記：\n% expansion:<category>"),
        ("Mode C — 加深", ACCENT,
         "已簽核章節再加深度。",
         "marker 必含：\n[pass: enrichment]"),
        ("Mode B — 審核", CORAL,
         "對每筆 marker 給判決。",
         "Keep / Rewrite / Move / Cut\nMove 只提案，不執行"),
    ]
    for i, (head, color, line1, line2) in enumerate(descs):
        x = Inches(0.55) + i * (col_w + gap)
        add_round_rect(s, x, desc_y, col_w, desc_h, fill=CARD_BG2, radius=0.04)
        add_rect(s, x, desc_y, col_w, Inches(0.08), fill=color)
        add_text(s, x + Inches(0.3), desc_y + Inches(0.25),
                 col_w - Inches(0.5), Inches(0.5),
                 [{"text": head, "size": 18, "color": color,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, x + Inches(0.3), desc_y + Inches(0.85),
                 col_w - Inches(0.5), Inches(0.5),
                 [{"text": line1, "size": 14, "color": BODY,
                   "font": FONT_BODY_CN}], line_spacing=1.3)
        add_text(s, x + Inches(0.3), desc_y + Inches(1.4),
                 col_w - Inches(0.5), Inches(0.85),
                 [{"text": line2, "size": 12.5, "color": SUBTLE,
                   "font": "Consolas"}], line_spacing=1.35)

    add_footer(s)


# ===========================================================
# Slide 5 — Expansion markers
# ===========================================================
def slide_expansion():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "% expansion: 標記 — 讓擴充可被審核")

    # 3x3 grid of category boxes
    cats = [
        ("history",     "數學史"),
        ("application", "應用連結"),
        ("formula",     "衍生公式"),
        ("summary",     "綜合段落"),
        ("figure",      "圖示"),
        ("example",     "補充例題"),
        ("intuition",   "直覺鋪陳"),
        ("strategy",    "解題策略框"),
        ("caution",     "警示框"),
    ]
    grid_y = Inches(1.5)
    cell_w = Inches(4.0)
    cell_h = Inches(0.9)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    start_x = Inches(0.55)
    for i, (en, cn) in enumerate(cats):
        col = i % 3
        row = i // 3
        x = start_x + col * (cell_w + gap_x)
        y = grid_y + row * (cell_h + gap_y)
        add_round_rect(s, x, y, cell_w, cell_h, fill=CARD_BG, radius=0.06)
        add_rect(s, x, y, Inches(0.08), cell_h, fill=ACCENT)
        add_text(s, x + Inches(0.3), y + Inches(0.16),
                 cell_w - Inches(0.4), Inches(0.4),
                 [{"text": en, "size": 17, "color": MATH_BLUE,
                   "bold": True, "font": "Consolas"}])
        add_text(s, x + Inches(0.3), y + Inches(0.5),
                 cell_w - Inches(0.4), Inches(0.4),
                 [{"text": cn, "size": 13, "color": SUBTLE,
                   "font": FONT_BODY_CN}])

    # Bottom: grep + philosophy
    bot_y = Inches(5.4)
    # grep card
    add_round_rect(s, Inches(0.55), bot_y, Inches(7.5), Inches(1.55),
                   fill=CARD_BG2, radius=0.04)
    add_text(s, Inches(0.85), bot_y + Inches(0.15),
             Inches(7.0), Inches(0.4),
             [{"text": "審查方式", "size": 13, "color": ACCENT,
               "bold": True, "font": FONT_BODY_CN}])
    add_text(s, Inches(0.85), bot_y + Inches(0.55),
             Inches(7.0), Inches(0.5),
             [{"text": 'grep "^% expansion:" chapters/chNN_*.tex',
               "size": 15, "color": MATH_BLUE, "font": "Consolas"}])
    add_text(s, Inches(0.85), bot_y + Inches(1.05),
             Inches(7.0), Inches(0.4),
             [{"text": "→ 一行掃出全部擴充，逐筆審視 Keep / Rewrite / Move / Cut",
               "size": 12, "color": SUBTLE, "font": FONT_BODY_CN}])

    # philosophy card
    add_round_rect(s, Inches(8.25), bot_y, Inches(4.55), Inches(1.55),
                   fill=CARD_BG2, radius=0.04, line=ACCENT)
    add_text(s, Inches(8.55), bot_y + Inches(0.15),
             Inches(4.0), Inches(0.4),
             [{"text": "擴充哲學", "size": 13, "color": ACCENT,
               "bold": True, "font": FONT_BODY_CN}])
    add_text(s, Inches(8.55), bot_y + Inches(0.55),
             Inches(4.0), Inches(0.95),
             [{"text": "寧可多寫並標記，",
               "size": 14, "color": BODY, "bold": True,
               "font": FONT_TITLE_CN}], line_spacing=1.3)
    add_text(s, Inches(8.55), bot_y + Inches(0.95),
             Inches(4.0), Inches(0.55),
             [{"text": "不要因怕越界而自我審查。",
               "size": 14, "color": BODY, "bold": True,
               "font": FONT_TITLE_CN}])

    add_footer(s)


# ===========================================================
# Slide 6 — Progress overview
# ===========================================================
def slide_progress():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "目前進度")

    # Top scope banner
    add_round_rect(s, Inches(0.55), Inches(1.45),
                   Inches(12.23), Inches(0.65),
                   fill=CARD_BG2, radius=0.06)
    add_text(s, Inches(0.85), Inches(1.45),
             Inches(11.6), Inches(0.65),
             [{"text": "範圍：Calc I + II + III  ·  共約 14 章",
               "size": 16, "color": BODY, "bold": True,
               "font": FONT_TITLE_CN}],
             anchor=MSO_ANCHOR.MIDDLE)

    # Done section (left) — 4 chapter cards stacked
    done_y = Inches(2.35)
    # header row
    add_text(s, Inches(0.55), done_y, Inches(6), Inches(0.4),
             [{"text": "✓  draft (4)", "size": 18, "color": EMERALD,
               "bold": True, "font": "Calibri"}])
    chs = [
        ("Ch 1", "Inverse Functions and Limits", "6 sections"),
        ("Ch 2", "Derivatives", "5 sections"),
        ("Ch 3", "Chain Rule & Trigonometric Derivatives", "3 sections"),
        ("Ch 4", "The Exponential and Logarithmic Functions", "5 sections"),
    ]
    cy = done_y + Inches(0.5)
    ch_h = Inches(0.85)
    ch_gap = Inches(0.12)
    for tag, title, meta in chs:
        add_round_rect(s, Inches(0.55), cy, Inches(6.15), ch_h,
                       fill=CARD_BG, radius=0.06)
        add_rect(s, Inches(0.55), cy, Inches(0.08), ch_h, fill=EMERALD)
        add_text(s, Inches(0.85), cy + Inches(0.1),
                 Inches(1.0), Inches(0.65),
                 [{"text": tag, "size": 16, "color": EMERALD,
                   "bold": True, "font": "Calibri"}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.85), cy + Inches(0.13),
                 Inches(4.7), Inches(0.4),
                 [{"text": title, "size": 14, "color": BODY,
                   "bold": True, "font": "Calibri"}])
        add_text(s, Inches(1.85), cy + Inches(0.5),
                 Inches(4.7), Inches(0.35),
                 [{"text": meta, "size": 11.5, "color": SUBTLE,
                   "font": "Calibri"}])
        cy += ch_h + ch_gap

    # TODO section (right)
    todo_y = done_y
    add_text(s, Inches(7.0), todo_y, Inches(6), Inches(0.4),
             [{"text": "○  TODO  Ch 5–14", "size": 18, "color": ACCENT,
               "bold": True, "font": "Calibri"}])
    todo_items = [
        ("Calc I 剩餘", "Applications of Differentiation · Integrals"),
        ("Calc II", "5 章（Applications, Techniques, ODEs, Parametric/Polar, Series）"),
        ("Calc III", "5 章（Vectors, Vector Functions, Partial Deriv., Multiple Int., Vector Calc.）"),
    ]
    ty = todo_y + Inches(0.5)
    th = Inches(1.18)
    tgap = Inches(0.12)
    for head, body in todo_items:
        add_round_rect(s, Inches(7.0), ty, Inches(5.78), th,
                       fill=CARD_BG2, radius=0.06)
        add_rect(s, Inches(7.0), ty, Inches(0.08), th, fill=ACCENT)
        add_text(s, Inches(7.3), ty + Inches(0.18),
                 Inches(5.3), Inches(0.4),
                 [{"text": head, "size": 14, "color": ACCENT,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, Inches(7.3), ty + Inches(0.55),
                 Inches(5.3), Inches(0.6),
                 [{"text": body, "size": 12, "color": BODY,
                   "font": "Calibri"}], line_spacing=1.35)
        ty += th + tgap

    add_text(s, Inches(7.0), ty + Inches(0.05),
             Inches(5.78), Inches(0.35),
             [{"text": "標題待前章穩定後再展開",
               "size": 11.5, "color": SUBTLE, "italic": True,
               "font": FONT_BODY_CN}])

    add_footer(s)


# ===========================================================
# Helper: chapter slide WITH image (for Ch1, Ch4)
# ===========================================================
def chapter_slide_with_image(idx, ch_num, title_en, sections, manuscript,
                              highlight, image_path, image_ratio,
                              image_caption, total_idx):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, f"Ch {ch_num} — {title_en}")

    # Status badge under title
    add_round_rect(s, Inches(0.85), Inches(1.4),
                   Inches(2.3), Inches(0.4),
                   fill=CARD_BG2, radius=0.20)
    add_text(s, Inches(0.85), Inches(1.4), Inches(2.3), Inches(0.4),
             [{"text": f"draft  ·  {len(sections)} sections",
               "size": 12, "color": EMERALD, "bold": True,
               "font": "Calibri"}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # LEFT: image — maximize within available content area (y=1.95 to ~6.95)
    img_y = Inches(1.95)
    img_h = Inches(5.00)                       # full content-area height
    img_w_in = round(5.00 * image_ratio, 3)    # portrait → narrower
    img_w = Inches(img_w_in)
    img_x = Inches(0.55)
    add_image_with_frame(s, image_path, img_x, img_y, img_w, img_h)

    # RIGHT: flowing text (no cramped cards). Vertical budget 2.05–6.85 = 4.80"
    right_x = img_x + img_w + Inches(0.35)
    right_w = SLIDE_W - right_x - Inches(0.55)

    cy = Inches(2.05)
    # Section heading
    add_text(s, right_x, cy, right_w, Inches(0.35),
             [{"text": "章節結構", "size": 14, "color": ACCENT,
               "bold": True, "font": FONT_TITLE_CN}])
    cy += Inches(0.40)
    # Sections list (compact rows)
    item_h = Inches(0.38)
    for code, name in sections:
        add_text(s, right_x + Inches(0.10), cy,
                 Inches(0.85), item_h,
                 [{"text": code, "size": 12, "color": MATH_BLUE,
                   "bold": True, "font": "Calibri"}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, right_x + Inches(1.0), cy,
                 right_w - Inches(1.0), item_h,
                 [{"text": name, "size": 12, "color": BODY,
                   "font": "Calibri"}],
                 anchor=MSO_ANCHOR.MIDDLE)
        cy += item_h

    cy += Inches(0.20)
    # Thin gold rule before manuscript info
    add_rect(s, right_x, cy, right_w, Inches(0.025), fill=ACCENT)
    cy += Inches(0.18)
    # Manuscript heading + body
    add_text(s, right_x, cy, right_w, Inches(0.30),
             [{"text": "Manuscript", "size": 12, "color": ACCENT,
               "bold": True, "font": "Calibri"}])
    cy += Inches(0.32)
    add_text(s, right_x, cy, right_w, Inches(0.40),
             [{"text": manuscript, "size": 11, "color": BODY,
               "font": FONT_BODY_CN}], line_spacing=1.3)
    cy += Inches(0.48)
    # Highlight heading + body
    add_text(s, right_x, cy, right_w, Inches(0.30),
             [{"text": "亮點", "size": 12, "color": TITLE,
               "bold": True, "font": FONT_TITLE_CN}])
    cy += Inches(0.32)
    add_text(s, right_x, cy, right_w, Inches(1.40),
             [{"text": highlight, "size": 11, "color": SUBTLE,
               "font": FONT_BODY_CN}], line_spacing=1.4)

    add_footer(s)


# ===========================================================
# Helper: chapter slide (text-only, for Ch2, Ch3)
# ===========================================================
def chapter_slide(idx, ch_num, title_en, sections, manuscript, highlight,
                  total_idx):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, f"Ch {ch_num} — {title_en}")

    # Status badge under title
    add_round_rect(s, Inches(0.85), Inches(1.4),
                   Inches(2.3), Inches(0.4),
                   fill=CARD_BG2, radius=0.20)
    add_text(s, Inches(0.85), Inches(1.4), Inches(2.3), Inches(0.4),
             [{"text": f"draft  ·  {len(sections)} sections",
               "size": 12, "color": EMERALD, "bold": True,
               "font": "Calibri"}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Left column: sections list
    left_x = Inches(0.55)
    left_w = Inches(7.5)
    sec_y = Inches(2.1)
    add_text(s, left_x + Inches(0.3), sec_y, left_w, Inches(0.4),
             [{"text": "章節結構", "size": 14, "color": ACCENT,
               "bold": True, "font": FONT_TITLE_CN}])
    list_y = sec_y + Inches(0.5)
    add_round_rect(s, left_x, list_y, left_w, Inches(4.4),
                   fill=CARD_BG, radius=0.04)
    inner_y = list_y + Inches(0.2)
    for code, name in sections:
        add_text(s, left_x + Inches(0.3), inner_y,
                 Inches(0.9), Inches(0.4),
                 [{"text": code, "size": 13, "color": MATH_BLUE,
                   "bold": True, "font": "Calibri"}])
        add_text(s, left_x + Inches(1.1), inner_y,
                 left_w - Inches(1.3), Inches(0.4),
                 [{"text": name, "size": 13, "color": BODY,
                   "font": "Calibri"}])
        inner_y += Inches(0.5)

    # Right column: manuscript + highlight
    right_x = Inches(8.3)
    right_w = Inches(4.48)
    # Manuscript card
    add_round_rect(s, right_x, sec_y, right_w, Inches(2.0),
                   fill=CARD_BG2, radius=0.04)
    add_rect(s, right_x, sec_y, right_w, Inches(0.06), fill=ACCENT)
    add_text(s, right_x + Inches(0.3), sec_y + Inches(0.2),
             right_w - Inches(0.5), Inches(0.4),
             [{"text": "Manuscript", "size": 13, "color": ACCENT,
               "bold": True, "font": "Calibri"}])
    add_text(s, right_x + Inches(0.3), sec_y + Inches(0.65),
             right_w - Inches(0.5), Inches(1.3),
             [{"text": manuscript, "size": 12, "color": BODY,
               "font": FONT_BODY_CN}], line_spacing=1.4)

    # Highlight card
    hl_y = sec_y + Inches(2.2)
    add_round_rect(s, right_x, hl_y, right_w, Inches(2.7),
                   fill=CARD_BG, radius=0.04, line=TITLE)
    add_rect(s, right_x, hl_y, right_w, Inches(0.06), fill=TITLE)
    add_text(s, right_x + Inches(0.3), hl_y + Inches(0.2),
             right_w - Inches(0.5), Inches(0.4),
             [{"text": "亮點與註記", "size": 13, "color": TITLE,
               "bold": True, "font": FONT_TITLE_CN}])
    add_text(s, right_x + Inches(0.3), hl_y + Inches(0.65),
             right_w - Inches(0.5), Inches(2.0),
             [{"text": highlight, "size": 12, "color": BODY,
               "font": FONT_BODY_CN}], line_spacing=1.4)

    add_footer(s)


# ===========================================================
# Slides 7–10 — chapters
# ===========================================================
def slide_ch1():
    chapter_slide(
        idx=7, ch_num=1, title_en="Inverse Functions and Limits",
        sections=[
            ("§1.1", "Inverse Functions and One-to-One Functions"),
            ("§1.2", "Inverse Trigonometric Functions"),
            ("§1.3", "Limits"),
            ("§1.4", "One-Sided and Infinite Limits"),
            ("§1.5", "Limit Laws"),
            ("§1.6", "The Precise Definition of a Limit"),
        ],
        manuscript="既存 LaTeX；以反向工程方式填入 roadmap。後續編修時，LaTeX 與 roadmap 同步更新。",
        highlight=(
            "・完整 ε–δ 處理\n"
            "・配套兩支 Manim 示範影片：§1.1 與 §1.6\n"
            "・作為其他章節的「register 校準器」"
        ),
        total_idx=7,
    )


def slide_ch2():
    chapter_slide(
        idx=8, ch_num=2, title_en="Derivatives",
        sections=[
            ("§2.1", "The Tangent Line and the Derivative at a Point"),
            ("§2.2", "The Derivative as a Function"),
            ("§2.3", "Differentiability, Continuity, and Higher Derivatives"),
            ("§2.4", "Derivatives of Polynomials and the Exponential Function"),
            ("§2.5", "The Product and Quotient Rules"),
        ],
        manuscript="2026-04-27 收到（13 頁手稿）。涵蓋切線動機 → 基本法則 → e^x 與乘除法則。",
        highlight=(
            "・三角 / 鏈鎖 / 隱函數 → 拆出至 Ch 3、Ch 4\n"
            "・e^x 採級數定義（與 Ch 4 對接）\n"
            "・higher derivatives 收進 §2.3 末段"
        ),
        total_idx=8,
    )


def slide_ch3():
    chapter_slide(
        idx=9, ch_num=3, title_en="Chain Rule & Trigonometric Derivatives",
        sections=[
            ("§3.1", "Derivatives of the Sine and Cosine Functions"),
            ("§3.2", "The Chain Rule"),
            ("§3.3", "Applications of the Chain Rule (ln x, x^x, arcsin x)"),
        ],
        manuscript="2023-10-28-chainRule。涵蓋 sin/cos 導數、鏈鎖法則證明、三類應用。",
        highlight=(
            "・squeeze theorem 用於 sin θ / θ → 1\n"
            "・餘式形式可微分定義（chain-rule 證明用）\n"
            "・ln x 暫以「e^x 反函數」處理；嚴謹版見 Ch 4"
        ),
        total_idx=9,
    )


def slide_ch4():
    chapter_slide(
        idx=10, ch_num=4, title_en="The Exponential and Logarithmic Functions",
        sections=[
            ("§4.1", "Construction of the Exponential Function"),
            ("§4.2", "Continuity and the Exponent Law for e^x"),
            ("§4.3", "The Derivative of e^x"),
            ("§4.4", "Rolle's Theorem and the Mean Value Theorem"),
            ("§4.5", "Monotonicity and the Logarithmic Function"),
        ],
        manuscript="2023-11-4-ExponentialFunction。級數定義 e^x → Cauchy ⇔ convergent → MVT → ln x。",
        highlight=(
            "・以 Bolzano–Weierstrass 證 Cauchy ⇔ 收斂\n"
            "・MVT 暫居 Ch 4，未來可能搬至「微分應用」章\n"
            "・閉合 Ch 2 / Ch 3 對 e^x、ln x 的非嚴謹引用"
        ),
        total_idx=10,
    )


# ===========================================================
# Slide 11 — Manim pipeline
# ===========================================================
def slide_manim():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "講義 → Manim 影片")

    # 5-step horizontal flow
    flow_y = Inches(1.55)
    step_w = Inches(2.20)
    step_h = Inches(1.05)
    gap = Inches(0.20)
    start_x = Inches(0.55)
    steps = [
        ("01", ".tex 完稿", "章節 LaTeX 定稿"),
        ("02", "手寫 storyboard", "YAML，不自動生"),
        ("03", "預覽單一 scene", "校稿 narration.md"),
        ("04", "TTS 合成", "Coqui XTTS / F5"),
        ("05", "最終 render", "1920×1080 / 30fps"),
    ]
    for i, (num, head, sub) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        add_round_rect(s, x, flow_y, step_w, step_h, fill=CARD_BG, radius=0.05)
        add_rect(s, x, flow_y, step_w, Inches(0.06), fill=TITLE)
        add_text(s, x + Inches(0.2), flow_y + Inches(0.15),
                 Inches(0.7), Inches(0.35),
                 [{"text": num, "size": 13, "color": ACCENT,
                   "bold": True, "font": "Calibri"}])
        add_text(s, x + Inches(0.2), flow_y + Inches(0.45),
                 step_w - Inches(0.4), Inches(0.35),
                 [{"text": head, "size": 14, "color": BODY,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, x + Inches(0.2), flow_y + Inches(0.78),
                 step_w - Inches(0.4), Inches(0.3),
                 [{"text": sub, "size": 11, "color": SUBTLE,
                   "font": FONT_BODY_CN}])
        # arrow between
        if i < len(steps) - 1:
            ax = x + step_w + Inches(0.04)
            ay = flow_y + step_h // 2
            add_arrow(s, ax, ay, ax + gap - Inches(0.08), ay,
                      color=ACCENT, width=2)

    # Templates list (left)
    tmpl_y = Inches(3.05)
    add_round_rect(s, Inches(0.55), tmpl_y, Inches(7.5), Inches(3.8),
                   fill=CARD_BG, radius=0.04)
    add_rect(s, Inches(0.55), tmpl_y, Inches(7.5), Inches(0.06), fill=TITLE)
    add_text(s, Inches(0.85), tmpl_y + Inches(0.2),
             Inches(7.0), Inches(0.4),
             [{"text": "9 個 Scene Templates", "size": 16, "color": TITLE,
               "bold": True, "font": FONT_TITLE_CN}])
    templates = [
        "title_bullets", "definition_math", "example_walkthrough",
        "graph_focus", "procedure_steps", "recap_cards",
        "section_transition", "theorem_proof", "comparison",
    ]
    grid_start_y = tmpl_y + Inches(0.7)
    cw = Inches(2.35)
    ch = Inches(0.45)
    cgap_x = Inches(0.08)
    cgap_y = Inches(0.18)
    for i, t in enumerate(templates):
        col = i % 3
        row = i // 3
        cx = Inches(0.85) + col * (cw + cgap_x)
        cy = grid_start_y + row * (ch + cgap_y)
        add_round_rect(s, cx, cy, cw, ch, fill=CARD_BG2, radius=0.10)
        add_text(s, cx, cy, cw, ch,
                 [{"text": t, "size": 12, "color": MATH_BLUE,
                   "bold": True, "font": "Consolas"}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Single source of truth note
    add_text(s, Inches(0.85), tmpl_y + Inches(3.15),
             Inches(7.0), Inches(0.4),
             [{"text": "單一真相：storyboard YAML",
               "size": 12, "color": SUBTLE, "italic": True,
               "font": FONT_BODY_CN}])

    # Visual style card (right)
    vis_y = Inches(3.05)
    add_round_rect(s, Inches(8.3), vis_y, Inches(4.48), Inches(3.8),
                   fill=CARD_BG2, radius=0.04, line=ACCENT)
    add_rect(s, Inches(8.3), vis_y, Inches(4.48), Inches(0.06), fill=ACCENT)
    add_text(s, Inches(8.6), vis_y + Inches(0.2),
             Inches(4.0), Inches(0.4),
             [{"text": "Midnight Canvas 視覺", "size": 16, "color": ACCENT,
               "bold": True, "font": FONT_TITLE_CN}])
    add_text(s, Inches(8.6), vis_y + Inches(0.7),
             Inches(4.0), Inches(0.45),
             [{"text": "深底  ·  發光數學  ·  零裝飾",
               "size": 13, "color": BODY, "font": FONT_TITLE_CN}])
    # Color semantics
    sem_y = vis_y + Inches(1.25)
    sems = [
        ("cyan", "定義 / 命題", TITLE),
        ("gold", "定理 / 關鍵結果", ACCENT),
        ("coral", "警示 / 反例", CORAL),
        ("emerald", "證畢 / 驗證", EMERALD),
    ]
    for i, (en, cn, col) in enumerate(sems):
        py = sem_y + i * Inches(0.55)
        # color swatch
        add_rect(s, Inches(8.6), py + Inches(0.08),
                 Inches(0.3), Inches(0.3), fill=col)
        add_text(s, Inches(9.05), py,
                 Inches(1.5), Inches(0.45),
                 [{"text": en, "size": 13, "color": col,
                   "bold": True, "font": "Calibri"}],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.4), py,
                 Inches(2.2), Inches(0.45),
                 [{"text": cn, "size": 13, "color": BODY,
                   "font": FONT_BODY_CN}],
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s)


# ===========================================================
# Slide 12 — Summary
# ===========================================================
def slide_summary():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "總結")

    items = [
        ("01", "三模式機制",
         "讓多人手稿能被 Claude 安全擴充又可審核：A 起草 / C 加深 / B 審核",
         TITLE),
        ("02", "目前進度",
         "Calc I 前 4 章 draft；Ch 5–14 待依序展開（TODO）",
         EMERALD),
        ("03", "雙產出共用一份 LaTeX",
         "講義（A4 PDF）與 Manim 影片同源；storyboard YAML 為影片真相",
         ACCENT),
        ("04", "GitHub + CI 五道把關",
         "lint × 3、storyboard lint、latexmk；避免格式 / 交叉引用 drift",
         CORAL),
    ]
    y = Inches(1.55)
    h = Inches(1.10)
    gap = Inches(0.18)
    for num, head, body, color in items:
        # number + colored bar
        add_round_rect(s, Inches(0.55), y, Inches(1.0), h,
                       fill=CARD_BG2, radius=0.10)
        add_text(s, Inches(0.55), y, Inches(1.0), h,
                 [{"text": num, "size": 32, "color": color,
                   "bold": True, "font": "Calibri"}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # row card
        add_round_rect(s, Inches(1.7), y, Inches(11.08), h,
                       fill=CARD_BG, radius=0.04)
        add_rect(s, Inches(1.7), y, Inches(0.08), h, fill=color)
        add_text(s, Inches(2.0), y + Inches(0.18),
                 Inches(10.7), Inches(0.5),
                 [{"text": head, "size": 19, "color": BODY,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, Inches(2.0), y + Inches(0.65),
                 Inches(10.7), Inches(0.45),
                 [{"text": body, "size": 13.5, "color": SUBTLE,
                   "font": FONT_BODY_CN}], line_spacing=1.3)
        y += h + gap

    # Closing tagline beneath last row, above footer rule.
    add_text(s, Inches(0.55), Inches(6.65), Inches(12.23), Inches(0.30),
             [{"text": "Thank you  ·  Q & A",
               "size": 13, "color": ACCENT, "italic": True,
               "font": "Calibri"}],
             align=PP_ALIGN.CENTER)

    add_footer(s)


# ===========================================================
# Slide  4 — Environment cheat-sheet showcase
# ===========================================================
def slide_env_showcase():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "規範化目標：spec 的 12 個語意環境")

    # Two-column layout: large image on the left, side-panel context on right
    img_x = Inches(0.55)
    img_w = Inches(8.5)
    img_h = Inches(5.52)            # 8.5 / 1.539 ≈ 5.52
    img_y = Inches(1.45)
    add_image_with_frame(s, IMG_ENV_TBL, img_x, img_y, img_w, img_h)
    add_text(s, img_x, img_y + img_h + Inches(0.10),
             img_w, Inches(0.32),
             [{"text": "節錄自 CONTENT_QUICKSTART  ·  專案規定的 12 個語意環境",
               "size": 11, "color": SUBTLE, "italic": True,
               "font": FONT_BODY_CN}])

    # Right side panel: brief context bullets
    rx = Inches(9.30)
    rw = Inches(3.50)
    add_text(s, rx, Inches(1.55), rw, Inches(0.4),
             [{"text": "為什麼要這張表？",
               "size": 14, "color": ACCENT, "bold": True,
               "font": FONT_TITLE_CN}])
    bullets = [
        ("一張表規定一致性",
         "Claude 把每段手稿映射到\n其中一個環境，不發明新的"),
        ("強制配對",
         "example + solution 必須包進\nworkedexample，不可單獨出現"),
        ("獨立計數",
         "definition / theorem / proposition\n各有自己的計數器，方便查找"),
        ("視覺語意",
         "caution / strategy 帶左色塊；\nremark 是真的「附註」而非主軸"),
    ]
    by = Inches(2.05)
    for head, body in bullets:
        add_round_rect(s, rx, by, rw, Inches(0.95),
                       fill=CARD_BG, radius=0.06)
        add_rect(s, rx, by, Inches(0.06), Inches(0.95), fill=ACCENT)
        add_text(s, rx + Inches(0.20), by + Inches(0.10),
                 rw - Inches(0.3), Inches(0.35),
                 [{"text": head, "size": 12, "color": BODY,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, rx + Inches(0.20), by + Inches(0.42),
                 rw - Inches(0.3), Inches(0.55),
                 [{"text": body, "size": 10.5, "color": SUBTLE,
                   "font": FONT_BODY_CN}], line_spacing=1.25)
        by += Inches(1.10)

    add_footer(s)


# ===========================================================
# Slide  8 — Handout showcase (Ch1 §1.6 + Ch4 章首 side by side)
# ===========================================================
def slide_handout_showcase():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "講義成書樣貌：實際印出來的樣子")

    # Two screenshots side by side, both portrait, h=5.10
    img_h_in = 5.10
    # Ch1 §1.6 (shot_03, ratio 637/810 = 0.786)
    w_left = round(img_h_in * (637 / 810), 3)
    # Ch4 章首 (shot_02, ratio 543/742 = 0.732)
    w_right = round(img_h_in * (543 / 742), 3)
    gap = 0.55
    total_w = w_left + gap + w_right
    start_x = (13.333 - total_w) / 2

    img_y = Inches(1.55)
    img_h = Inches(img_h_in)

    x_left = Inches(round(start_x, 3))
    x_right = Inches(round(start_x + w_left + gap, 3))
    add_image_with_frame(s, IMG_CH1_PAGE,
                         x_left, img_y, Inches(w_left), img_h)
    add_image_with_frame(s, IMG_CH4_OPEN,
                         x_right, img_y, Inches(w_right), img_h)

    # Captions under each image
    cap_y = img_y + img_h + Inches(0.10)
    add_text(s, x_left, cap_y, Inches(w_left), Inches(0.3),
             [{"text": "Ch 1 — §1.6  ε–δ strategy + worked example",
               "size": 12, "color": TITLE, "bold": True,
               "font": "Calibri"}],
             align=PP_ALIGN.CENTER)
    add_text(s, x_right, cap_y, Inches(w_right), Inches(0.3),
             [{"text": "Ch 4 — 章首與 learning outcomes",
               "size": 12, "color": TITLE, "bold": True,
               "font": "Calibri"}],
             align=PP_ALIGN.CENTER)

    # Bottom tagline
    add_text(s, Inches(0.55), cap_y + Inches(0.50),
             Inches(12.23), Inches(0.32),
             [{"text": "Stewart / Rogawski 風格 ‧ 直接由 LaTeX 編譯成 A4 PDF ‧ 章首 / strategy / workedexample 等環境一致呈現",
               "size": 11, "color": SUBTLE, "italic": True,
               "font": FONT_BODY_CN}],
             align=PP_ALIGN.CENTER)

    add_footer(s)


# ===========================================================
# Slide 14 — GitHub narration.md showcase
# ===========================================================
def slide_github_showcase():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_title(s, "narration.md：跨 Markdown 校稿，再寫回 YAML")

    # Wide image (shot_07 ratio ≈ 2.014)
    img_h = Inches(4.60)
    img_w_in = round(4.60 * (1915 / 951), 3)   # ≈ 9.27
    img_w = Inches(img_w_in)
    img_x = Inches(round((13.333 - img_w_in) / 2, 3))
    img_y = Inches(1.45)
    add_image_with_frame(s, IMG_GITHUB, img_x, img_y, img_w, img_h)

    # Three-step workflow row beneath the image
    steps = [
        ("01", "渲染時自動匯出",
         "manim_render_lesson.py --with-audio\n會把 narration 匯出到 .md"),
        ("02", "在 Markdown 直接校稿",
         "字句、語氣、發音改起來\n比改 YAML 直觀很多"),
        ("03", "同步寫回 storyboard",
         "manim_sync_narration_back.py\n把改好的內容寫回 YAML"),
    ]
    sy = img_y + img_h + Inches(0.30)
    sw = Inches(4.05)
    sgap = Inches(0.18)
    sx0 = Inches(round((13.333 - 3 * 4.05 - 2 * 0.18) / 2, 3))
    for i, (num, head, body) in enumerate(steps):
        sx = sx0 + i * (sw + sgap)
        add_round_rect(s, sx, sy, sw, Inches(0.95),
                       fill=CARD_BG, radius=0.06)
        add_rect(s, sx, sy, Inches(0.08), Inches(0.95), fill=ACCENT)
        add_text(s, sx + Inches(0.25), sy + Inches(0.10),
                 Inches(0.6), Inches(0.4),
                 [{"text": num, "size": 14, "color": ACCENT,
                   "bold": True, "font": "Calibri"}])
        add_text(s, sx + Inches(0.85), sy + Inches(0.10),
                 sw - Inches(1.0), Inches(0.4),
                 [{"text": head, "size": 13, "color": BODY,
                   "bold": True, "font": FONT_TITLE_CN}])
        add_text(s, sx + Inches(0.25), sy + Inches(0.50),
                 sw - Inches(0.4), Inches(0.45),
                 [{"text": body, "size": 10.5, "color": SUBTLE,
                   "font": "Consolas"}], line_spacing=1.25)

    add_footer(s)


# ---------- Build all ----------
slide_cover()                # 1
slide_overview()             # 2
slide_manuscript()           # 3
slide_env_showcase()         # 4   NEW
slide_three_modes()          # 5
slide_expansion()            # 6
slide_progress()             # 7
slide_handout_showcase()     # 8   NEW
slide_ch1()                  # 9
slide_ch2()                  # 10
slide_ch3()                  # 11
slide_ch4()                  # 12
slide_manim()                # 13
slide_github_showcase()      # 14  NEW
slide_summary()              # 15

import sys
out = r"C:\Users\user\Desktop\Calculus_handout\calculus_handout_workflow.pptx"
prs.save(out)
print("Saved:", out)
